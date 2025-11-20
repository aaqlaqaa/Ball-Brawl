from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add a title slide with gradient background"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background shape
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 45
    fill.gradient_stops[0].color.rgb = RGBColor(41, 128, 185)  # Blue
    fill.gradient_stops[1].color.rgb = RGBColor(109, 213, 250)  # Light blue
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.2),
        Inches(8), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_points, bg_color1, bg_color2):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add gradient background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135
    fill.gradient_stops[0].color.rgb = RGBColor(*bg_color1)
    fill.gradient_stops[1].color.rgb = RGBColor(*bg_color2)
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(9), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add content box
    content_box = slide.shapes.add_textbox(
        Inches(1), Inches(2),
        Inches(8), Inches(5)
    )
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = point
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.level = 0
        p.space_before = Pt(12)
    
    return slide

def add_section_slide(prs, section_title, bg_color1, bg_color2):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Add gradient background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 90
    fill.gradient_stops[0].color.rgb = RGBColor(*bg_color1)
    fill.gradient_stops[1].color.rgb = RGBColor(*bg_color2)
    
    # Add section title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(3),
        Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = section_title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    return slide

# Slide 1: Title Slide
add_title_slide(prs, "Marketing Fundamentals", "Building Blocks of Successful Marketing")

# Slide 2: What is Marketing?
add_section_slide(prs, "What is Marketing?", (231, 76, 60), (192, 57, 43))

add_content_slide(prs, "Definition of Marketing", [
    "• The process of creating, communicating, and delivering value to customers",
    "• Building and maintaining profitable customer relationships",
    "• Understanding customer needs and wants",
    "• Creating solutions that satisfy those needs",
    "• More than just advertising - it's a complete business philosophy"
], (52, 73, 94), (44, 62, 80))

# Slide 3: The 4 Ps of Marketing
add_section_slide(prs, "The 4 Ps of Marketing", (46, 204, 113), (39, 174, 96))

add_content_slide(prs, "Product", [
    "• What you're selling - goods, services, or experiences",
    "• Features, quality, design, and branding",
    "• Product lifecycle management",
    "• Differentiation from competitors",
    "• Meeting customer needs and solving problems"
], (155, 89, 182), (142, 68, 173))

add_content_slide(prs, "Price", [
    "• What customers pay for your product or service",
    "• Pricing strategies: premium, penetration, competitive",
    "• Perceived value vs. actual cost",
    "• Discounts, payment terms, and financing options",
    "• Impact on brand positioning and profitability"
], (52, 152, 219), (41, 128, 185))

add_content_slide(prs, "Place (Distribution)", [
    "• Where and how customers can access your product",
    "• Distribution channels: direct, retail, online, wholesale",
    "• Logistics and supply chain management",
    "• Market coverage and accessibility",
    "• Omnichannel strategies for modern consumers"
], (230, 126, 34), (211, 84, 0))

add_content_slide(prs, "Promotion", [
    "• How you communicate with your target audience",
    "• Advertising, public relations, sales promotions",
    "• Digital marketing: social media, email, content marketing",
    "• Personal selling and direct marketing",
    "• Integrated marketing communications"
], (231, 76, 60), (192, 57, 43))

# Slide 4: Market Segmentation
add_section_slide(prs, "Market Segmentation", (241, 196, 15), (243, 156, 18))

add_content_slide(prs, "Understanding Your Audience", [
    "• Demographic: age, gender, income, education",
    "• Geographic: location, climate, urban vs. rural",
    "• Psychographic: lifestyle, values, personality, interests",
    "• Behavioral: purchase patterns, brand loyalty, usage rate",
    "• Targeting the right segments for maximum impact"
], (26, 188, 156), (22, 160, 133))

# Slide 5: Consumer Behavior
add_section_slide(prs, "Consumer Behavior", (155, 89, 182), (142, 68, 173))

add_content_slide(prs, "Understanding the Buyer's Journey", [
    "• Awareness: recognizing a need or problem",
    "• Consideration: researching and evaluating options",
    "• Decision: making the purchase choice",
    "• Post-purchase: evaluation and loyalty building",
    "• Factors influencing decisions: cultural, social, personal, psychological"
], (52, 73, 94), (44, 62, 80))

# Slide 6: Brand Management
add_section_slide(prs, "Brand Management", (230, 126, 34), (211, 84, 0))

add_content_slide(prs, "Building a Strong Brand", [
    "• Brand identity: name, logo, visual elements, voice",
    "• Brand positioning: unique place in customer's mind",
    "• Brand equity: value derived from consumer perception",
    "• Consistency across all touchpoints",
    "• Emotional connections and brand loyalty"
], (231, 76, 60), (192, 57, 43))

# Slide 7: Digital Marketing
add_section_slide(prs, "Digital Marketing", (52, 152, 219), (41, 128, 185))

add_content_slide(prs, "Marketing in the Digital Age", [
    "• Social media marketing: engagement and community building",
    "• Content marketing: valuable, relevant content creation",
    "• SEO & SEM: search engine visibility and paid advertising",
    "• Email marketing: personalized communication at scale",
    "• Analytics and data-driven decision making"
], (46, 204, 113), (39, 174, 96))

# Slide 8: Marketing Metrics
add_section_slide(prs, "Measuring Success", (155, 89, 182), (142, 68, 173))

add_content_slide(prs, "Key Performance Indicators (KPIs)", [
    "• ROI (Return on Investment): revenue vs. marketing spend",
    "• Customer Acquisition Cost (CAC)",
    "• Customer Lifetime Value (CLV)",
    "• Conversion rates and engagement metrics",
    "• Brand awareness and market share"
], (52, 73, 94), (44, 62, 80))

# Slide 9: Marketing Strategy
add_section_slide(prs, "Creating a Marketing Strategy", (26, 188, 156), (22, 160, 133))

add_content_slide(prs, "Strategic Planning Process", [
    "• Situation analysis: SWOT (Strengths, Weaknesses, Opportunities, Threats)",
    "• Setting SMART goals: Specific, Measurable, Achievable, Relevant, Time-bound",
    "• Target market selection and positioning",
    "• Marketing mix development (4 Ps)",
    "• Implementation, monitoring, and adjustment"
], (230, 126, 34), (211, 84, 0))

# Slide 10: Closing Slide
add_section_slide(prs, "Thank You!", (41, 128, 185), (109, 213, 250))

add_content_slide(prs, "Key Takeaways", [
    "• Marketing is about creating value for customers",
    "• Master the 4 Ps: Product, Price, Place, Promotion",
    "• Understand your audience through segmentation",
    "• Build strong brands that resonate emotionally",
    "• Embrace digital channels and data-driven decisions",
    "• Measure, analyze, and continuously improve"
], (46, 204, 113), (39, 174, 96))

# Save presentation
prs.save('/vercel/sandbox/Marketing_Fundamentals.pptx')
print("✅ PowerPoint presentation created successfully!")
print("📁 File saved as: Marketing_Fundamentals.pptx")
print("📊 Total slides: 17")
print("🎨 Features: Gradient backgrounds, professional color schemes, structured content")
